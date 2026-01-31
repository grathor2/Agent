"""Generate test data files for the agent system."""
import os
from pathlib import Path
from faker import Faker
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.units import inch
import docx
from pptx import Presentation
from pptx.util import Inches, Pt
import random

fake = Faker()

# Output directories
OUTPUT_DIR = Path(__file__).parent.parent / "data" / "generated"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

def generate_pdf(filename: str, num_pages: int = 6):
    """Generate a PDF file with multiple pages."""
    filepath = OUTPUT_DIR / filename
    c = canvas.Canvas(str(filepath), pagesize=letter)
    width, height = letter
    
    for page_num in range(num_pages):
        # Title
        c.setFont("Helvetica-Bold", 16)
        c.drawString(1*inch, height - 1*inch, f"Document Page {page_num + 1}")
        
        # Content
        c.setFont("Helvetica", 12)
        y_position = height - 1.5*inch
        
        for i in range(20):
            text = fake.paragraph(nb_sentences=3)
            text_lines = text.split('\n')
            for line in text_lines[:3]:
                if y_position < 1*inch:
                    break
                c.drawString(1*inch, y_position, line[:80])
                y_position -= 0.3*inch
        
        c.showPage()
    
    c.save()
    print(f"Generated PDF: {filename}")

def generate_docx(filename: str, num_pages: int = 6):
    """Generate a Word document."""
    filepath = OUTPUT_DIR / filename
    doc = docx.Document()
    
    for page_num in range(num_pages):
        # Title
        doc.add_heading(f'Document Section {page_num + 1}', level=1)
        
        # Content paragraphs
        for _ in range(5):
            doc.add_paragraph(fake.paragraph(nb_sentences=5))
        
        # Add page break (except for last page)
        if page_num < num_pages - 1:
            doc.add_page_break()
    
    doc.save(str(filepath))
    print(f"Generated DOCX: {filename}")

def generate_txt(filename: str, num_pages: int = 6):
    """Generate a text file."""
    filepath = OUTPUT_DIR / filename
    
    with open(filepath, 'w', encoding='utf-8') as f:
        for page_num in range(num_pages):
            f.write(f"=== Page {page_num + 1} ===\n\n")
            for _ in range(10):
                f.write(fake.paragraph(nb_sentences=4) + "\n\n")
    
    print(f"Generated TXT: {filename}")

def generate_pptx(filename: str, num_slides: int = 6):
    """Generate a PowerPoint presentation."""
    filepath = OUTPUT_DIR / filename
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    for slide_num in range(num_slides):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = f"Slide {slide_num + 1}: {fake.sentence()}"
        tf = content.text_frame
        tf.text = fake.paragraph(nb_sentences=3)
        
        # Add bullet points
        for _ in range(3):
            p = tf.add_paragraph()
            p.text = fake.sentence()
            p.level = 0
    
    prs.save(str(filepath))
    print(f"Generated PPTX: {filename}")

def generate_all_test_data():
    """Generate all test data files."""
    print("Generating test data files...")
    
    # Generate 30 PDFs
    for i in range(30):
        generate_pdf(f"test_document_{i+1}.pdf", num_pages=random.randint(5, 7))
    
    # Generate 30 DOCX files
    for i in range(30):
        generate_docx(f"test_document_{i+1}.docx", num_pages=random.randint(5, 7))
    
    # Generate 25 TXT files
    for i in range(25):
        generate_txt(f"test_document_{i+1}.txt", num_pages=random.randint(5, 7))
    
    # Generate 15 PPTX files
    for i in range(15):
        generate_pptx(f"test_presentation_{i+1}.pptx", num_slides=random.randint(5, 7))
    
    print(f"\n✅ Generated 100 test files in {OUTPUT_DIR}")
    print(f"   - 30 PDF files")
    print(f"   - 30 DOCX files")
    print(f"   - 25 TXT files")
    print(f"   - 15 PPTX files")

if __name__ == "__main__":
    generate_all_test_data()
