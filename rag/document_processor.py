"""Document processing with multi-format support including images."""
import os
from typing import List
try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
except ImportError:
    from langchain.text_splitter import RecursiveCharacterTextSplitter

try:
    from langchain_core.documents import Document
except ImportError:
    from langchain.schema import Document
try:
    import pytesseract
    HAS_OCR = True
except ImportError:
    HAS_OCR = False
    pytesseract = None

try:
    from PIL import Image
    from pdf2image import convert_from_path
    HAS_IMAGE_PROCESSING = True
except ImportError:
    HAS_IMAGE_PROCESSING = False

try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

import pypdf
from utils.logger import get_logger

logger = get_logger(__name__)

class DocumentProcessor:
    """Processes various document formats and extracts text with image OCR."""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        logger.info("DocumentProcessor initialized", 
                   chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    
    def extract_text_from_image(self, image_path: str) -> str:
        """Extract text from image using OCR."""
        if not HAS_OCR:
            logger.warning("OCR not available (pytesseract not installed)", image_path=image_path)
            return ""
        try:
            image = Image.open(image_path)
            text = pytesseract.image_to_string(image)
            logger.info("OCR completed", image_path=image_path, text_length=len(text))
            return text
        except Exception as e:
            logger.error("OCR failed", image_path=image_path, error=str(e))
            return ""
    
    def extract_images_from_pdf(self, pdf_path: str) -> List[str]:
        """Extract images from PDF and perform OCR."""
        if not HAS_IMAGE_PROCESSING:
            logger.warning("Image processing not available", pdf_path=pdf_path)
            return []
        image_texts = []
        try:
            images = convert_from_path(pdf_path)
            for i, image in enumerate(images):
                # Save temporary image
                temp_path = f"/tmp/pdf_image_{i}.png"
                image.save(temp_path, "PNG")
                # Extract text from image
                text = self.extract_text_from_image(temp_path)
                if text.strip():
                    image_texts.append(f"[Image {i+1} from PDF]: {text}")
                os.remove(temp_path)
            logger.info("PDF images processed", pdf_path=pdf_path, 
                       image_count=len(image_texts))
        except Exception as e:
            logger.error("PDF image extraction failed", pdf_path=pdf_path, error=str(e))
        return image_texts
    
    def process_pdf(self, file_path: str) -> List[Document]:
        """Process PDF file with text and image extraction."""
        documents = []
        try:
            # Extract text from PDF
            with open(file_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                text_content = []
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    text_content.append(f"Page {page_num + 1}:\n{text}")
                
                full_text = "\n\n".join(text_content)
            
            # Extract images from PDF
            image_texts = self.extract_images_from_pdf(file_path)
            if image_texts:
                full_text += "\n\n" + "\n\n".join(image_texts)
            
            # Create document
            doc = Document(
                page_content=full_text,
                metadata={"source": file_path, "type": "pdf", "pages": len(pdf_reader.pages)}
            )
            documents.append(doc)
            logger.info("PDF processed", file_path=file_path, 
                       text_length=len(full_text))
        except Exception as e:
            logger.error("PDF processing failed", file_path=file_path, error=str(e))
        return documents
    
    def process_docx(self, file_path: str) -> List[Document]:
        """Process Word document."""
        if not HAS_DOCX:
            logger.warning("DOCX processing not available (python-docx not installed)", file_path=file_path)
            return []
        documents = []
        try:
            doc = docx.Document(file_path)
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            full_text = "\n".join(paragraphs)
            
            # Extract images (embedded images in docx)
            # Note: python-docx doesn't directly extract images, but we can note their presence
            image_count = len(doc.inline_shapes)
            if image_count > 0:
                full_text += f"\n\n[Note: Document contains {image_count} embedded image(s)]"
            
            doc_obj = Document(
                page_content=full_text,
                metadata={"source": file_path, "type": "docx"}
            )
            documents.append(doc_obj)
            logger.info("DOCX processed", file_path=file_path, 
                       text_length=len(full_text))
        except Exception as e:
            logger.error("DOCX processing failed", file_path=file_path, error=str(e))
        return documents
    
    def process_txt(self, file_path: str) -> List[Document]:
        """Process plain text file."""
        documents = []
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            doc = Document(
                page_content=content,
                metadata={"source": file_path, "type": "txt"}
            )
            documents.append(doc)
            logger.info("TXT processed", file_path=file_path, 
                       text_length=len(content))
        except Exception as e:
            logger.error("TXT processing failed", file_path=file_path, error=str(e))
        return documents
    
    def process_pptx(self, file_path: str) -> List[Document]:
        """Process PowerPoint presentation."""
        if not HAS_PPTX:
            logger.warning("PPTX processing not available (python-pptx not installed)", file_path=file_path)
            return []
        documents = []
        try:
            prs = Presentation(file_path)
            slides_content = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"Slide {slide_num + 1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                slides_content.append(slide_text)
            
            full_text = "\n\n".join(slides_content)
            
            doc = Document(
                page_content=full_text,
                metadata={"source": file_path, "type": "pptx", 
                         "slides": len(prs.slides)}
            )
            documents.append(doc)
            logger.info("PPTX processed", file_path=file_path, 
                       text_length=len(full_text))
        except Exception as e:
            logger.error("PPTX processing failed", file_path=file_path, error=str(e))
        return documents
    
    def process_image(self, file_path: str) -> List[Document]:
        """Process image file with OCR."""
        documents = []
        try:
            text = self.extract_text_from_image(file_path)
            doc = Document(
                page_content=text if text else "[Image file - no text extracted]",
                metadata={"source": file_path, "type": "image"}
            )
            documents.append(doc)
            logger.info("Image processed", file_path=file_path, 
                       text_length=len(text))
        except Exception as e:
            logger.error("Image processing failed", file_path=file_path, error=str(e))
        return documents
    
    def process_file(self, file_path: str) -> List[Document]:
        """Process file based on extension."""
        if not file_path or '.' not in file_path:
            logger.warning("Invalid file path", file_path=file_path)
            return []
        
        extension = '.' + file_path.lower().split('.')[-1]
        
        processors = {
            ".pdf": self.process_pdf,
            ".docx": self.process_docx,
            ".txt": self.process_txt,
            ".pptx": self.process_pptx,
            ".png": self.process_image,
            ".jpg": self.process_image,
            ".jpeg": self.process_image,
        }
        
        processor = processors.get(extension)
        if not processor:
            logger.warning("Unsupported file type", file_path=file_path, extension=extension)
            return []
        
        return processor(file_path)
    
    def chunk_documents(self, documents: List[Document]) -> List[Document]:
        """Chunk documents with overlap strategy."""
        chunked_docs = []
        for doc in documents:
            chunks = self.text_splitter.split_documents([doc])
            chunked_docs.extend(chunks)
        
        logger.info("Documents chunked", 
                   original_count=len(documents), 
                   chunked_count=len(chunked_docs))
        return chunked_docs
